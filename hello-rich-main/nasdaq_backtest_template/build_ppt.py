import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from pathlib import Path

DATA_DIR = Path("data")
REPORT_DIR = Path("reports")

def build_ppt():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "纳斯达克100增强策略实盘回测报告（2015–2025）"
    slide.placeholders[1].text = f"作者：Eden Chen\n生成日期：{date.today()}"

    nav = pd.read_csv(REPORT_DIR / "backtest_nav.csv")
    plt.figure(figsize=(8,4))
    plt.plot(nav.index, nav.iloc[:,1])
    plt.title("策略累计净值")
    plt.tight_layout()
    plt.savefig(REPORT_DIR / "nav.png")
    plt.close()

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "累计净值曲线"
    slide.shapes.add_picture(str(REPORT_DIR / "nav.png"), Inches(0.5), Inches(1.5), width=Inches(8))

    ppt_path = REPORT_DIR / "NASDAQ100_Enhanced_Backtest_2015_2025.pptx"
    prs.save(ppt_path)
    print(f"✅ PPT saved to {ppt_path}")

if __name__ == "__main__":
    build_ppt()
